import os
import io
from textwrap import wrap
import streamlit as st
import pandas as pd
from datetime import datetime
from fpdf import FPDF
from langchain.schema import HumanMessage, AIMessage
from backend import build_or_load_index, query_faiss_index, get_llm_response

try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None

# ——— Custom PDF class ———
class ReportPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 16)
        self.cell(0, 10, 'Chat Conversation Report', ln=True, align='C')
        self.set_font('Arial', '', 10)
        self.cell(0, 8, datetime.now().strftime('%B %d, %Y %H:%M'),
                  ln=True, align='C')
        self.ln(5)
        self.set_draw_color(50, 50, 50)
        self.set_line_width(0.5)
        self.line(self.l_margin,
                  self.get_y(),
                  self.w - self.r_margin,
                  self.get_y())
        self.ln(5)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.set_text_color(100)
        self.cell(0, 10, f'Page {self.page_no()}', align='C')

def make_pdf_bytes(chat_history: list[HumanMessage|AIMessage]) -> bytes:
    pdf = ReportPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", "", 12)

    # table header
    ROLE_W = 40
    MSG_W = pdf.w - pdf.l_margin - pdf.r_margin - ROLE_W
    pdf.set_fill_color(200, 200, 200)
    pdf.set_font("Arial", "B", 12)
    pdf.cell(ROLE_W, 10, "Role", border=1, fill=True, align="C")
    pdf.cell(MSG_W, 10, "Message", border=1, fill=True, align="C", ln=True)

    fill = False
    for msg in chat_history:
        role = "User" if isinstance(msg, HumanMessage) else "Assistant"
        text = msg.content.replace("\n", " ")
        # ~2 chars per mm: adjust wrap width as needed
        lines = wrap(text, width=int(MSG_W / 2))

        pdf.set_font("Arial", "", 12)
        pdf.set_fill_color(245, 245, 245 if fill else 255)
        pdf.cell(ROLE_W, 10, role, border=1, fill=fill)
        for i, line in enumerate(lines):
            if i == 0:
                pdf.multi_cell(MSG_W, 10, line, border=1, fill=fill)
            else:
                pdf.cell(ROLE_W, 10, "", border=0)
                pdf.multi_cell(MSG_W, 10, line, border=1, fill=fill)
        fill = not fill

    raw = pdf.output(dest="S")
    return bytes(raw) if isinstance(raw, bytearray) else raw.encode("latin-1")


# ——— Streamlit setup ———
st.set_page_config(page_title="📘 Document Chat Agent", layout="wide")

# Load API key
try:
    GROQ_API_KEY = st.secrets["GROQ"]["API_KEY"]
except KeyError:
    st.error("🚫 GROQ API key not found. Add it under [GROQ] API_KEY in secrets.")
    st.stop()

# Sidebar
with st.sidebar:
    st.markdown("## ⚙️ Settings")
    uploaded_file = st.file_uploader(
        "📄 Upload document (PDF, TXT, DOCX, PPTX, CSV)",
        type=["pdf", "txt", "docx", "pptx", "csv"]
    )
    show_context = st.checkbox("🔎 Show Retrieved Contexts")
    if st.button("🧹 Clear Conversation"):
        st.session_state.pop("chat_history", None)

# Main title
st.markdown(
    "<h1 style='text-align:center;color:#4A90E2;'>💬 Conversational Document Chat Agent</h1>",
    unsafe_allow_html=True
)
st.markdown(
    "<p style='text-align:center;'>Upload a document and chat with its contents via a RAG pipeline.</p>",
    unsafe_allow_html=True
)
st.markdown("---")

# helpers
def extract_text(path, ftype):
    if ftype == "txt":
        return open(path, "r", encoding="utf-8").read()
    if ftype == "csv":
        return pd.read_csv(path).to_string()
    if ftype == "docx" and docx:
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if ftype == "pptx" and pptx:
        prs = pptx.Presentation(path)
        return "\n".join(
            shp.text for sl in prs.slides for shp in sl.shapes if hasattr(shp, "text")
        )
    return ""

def display_contexts(ctxs):
    st.markdown("### 🔍 Retrieved Contexts")
    for i, c in enumerate(ctxs, 1):
        st.markdown(
            f"<div style='background:#f0f2f6;padding:10px;"
            f"border-radius:8px;margin:4px 0;'>"
            f"<strong>Context {i}:</strong> {c}</div>",
            unsafe_allow_html=True
        )

# upload & chat logic
if uploaded_file:
    os.makedirs("docs", exist_ok=True)
    path = os.path.join("docs", uploaded_file.name)
    with open(path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    ext = uploaded_file.name.rsplit(".", 1)[1].lower()

    @st.cache_resource
    def load_index(p):
        return build_or_load_index(p)

    if ext in ["txt", "csv", "docx", "pptx"]:
        raw = extract_text(path, ext)
        txt_path = os.path.splitext(path)[0] + ".txt"
        with open(txt_path, "w", encoding="utf-8") as tf:
            tf.write(raw)
        index_path = txt_path
    else:
        index_path = path

    embedder, texts, index = load_index(index_path)
    st.session_state.setdefault("chat_history", [])

    # replay
    for m in st.session_state.chat_history:
        role = "user" if isinstance(m, HumanMessage) else "assistant"
        with st.chat_message(role):
            st.markdown(m.content)

    # input
    if q := st.chat_input("💬 Ask something..."):
        st.session_state.chat_history.append(HumanMessage(content=q))
        with st.chat_message("user"):
            st.markdown(q)

        ctxs = query_faiss_index(q, embedder, index, texts)
        if show_context:
            display_contexts(ctxs)

        resp = get_llm_response(
            api_key=GROQ_API_KEY,
            query=q,
            contexts=ctxs,
            chat_history=st.session_state.chat_history
        )
        st.session_state.chat_history.append(AIMessage(content=resp))
        with st.chat_message("assistant"):
            st.markdown(resp)

    # show download button once chat exists
    if st.session_state.chat_history:
        pdf_bytes = make_pdf_bytes(st.session_state.chat_history)
        st.sidebar.download_button(
            label="📥 Download Chat as PDF",
            data=pdf_bytes,
            file_name="chat_conversation_report.pdf",
            mime="application/pdf"
        )

else:
    st.info("📂 Upload a document to begin.")
